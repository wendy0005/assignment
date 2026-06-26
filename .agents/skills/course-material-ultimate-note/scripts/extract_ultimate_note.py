"""Extract text from PDFs, PPTXs, and binary PPTs in a folder and compile into an "ultimate note" markdown file."""
import fitz
import pptx
import olefile
import os
import glob
import re
import argparse


def extract_ppt_clean(filepath):
    """Extract text from old binary .ppt format via UTF-16LE string scanning."""
    try:
        ole = olefile.OleFileIO(filepath)
        if not ole.exists("PowerPoint Document"):
            ole.close()
            return None
        data = ole.openstream("PowerPoint Document").read()
        ole.close()
    except Exception:
        return None

    i = 0
    segments = []
    while i < len(data) - 1:
        if data[i + 1] == 0 and 0x20 <= data[i] <= 0x7E:
            start = i
            while i < len(data) - 1 and data[i + 1] == 0 and 0x20 <= data[i] <= 0x7E:
                i += 2
            text = data[start:i].decode("utf-16-le").strip()
            if len(text) >= 4:
                segments.append(text)
        else:
            i += 1

    exclude = {
        "rectangle", "footer placeholder", "slide number placeholder", "title",
        "content placeholder", "group", "oval", "picture", "textbox",
        "date placeholder", "notes placeholder", "slide image placeholder",
        "text placeholder", "arial", "times new roman", "wingdings",
        "rockwell", "lexend deca", "open sans", "ms pgothic", "wood type",
        "default design", "title placeholder",
    }

    content = []
    for s in segments:
        s_lower = s.lower().strip()
        if any(s_lower.startswith(e) for e in exclude):
            continue
        if s.startswith("___PPT") or s.startswith("__WPP"):
            continue
        if s.startswith("C:\\"):
            continue
        words = s.split()
        if len(words) >= 2 and any(c.islower() for c in s):
            alpha = sum(1 for c in s if c.isalpha()) / max(len(s), 1)
            if alpha > 0.4:
                content.append(s)

    return "\n".join(content) if content else None


def extract_pdf_text(filepath):
    """Extract text from a PDF file using PyMuPDF."""
    lines = []
    try:
        doc = fitz.open(filepath)
        for pn in range(len(doc)):
            text = doc[pn].get_text().strip()
            page_lines = [
                l.strip() for l in text.split("\n") if len(l.strip()) > 3
            ]
            if page_lines:
                lines.append(f"### Page {pn + 1}\n")
                lines.extend(page_lines)
                lines.append("")
        doc.close()
    except Exception as e:
        lines.append(f"*Error reading PDF: {e}*")
    return "\n".join(lines)


def extract_pptx_text(filepath):
    """Extract text from a PPTX file using python-pptx."""
    lines = []
    try:
        prs = pptx.Presentation(filepath)
        for sn, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                lines.append(f"### Slide {sn}\n")
                for t in texts:
                    lines.append(t)
                    lines.append("")
    except Exception as e:
        lines.append(f"*Error reading PPTX: {e}*")
    return "\n".join(lines)


def compile_ultimate_note(folder_path, folder_name=None):
    """
    Scan a folder for PDFs, PPTXs, PPTs and compile into '{folder_name} ultimate note.md'.

    Args:
        folder_path: Absolute path to the module folder
        folder_name: Display name (defaults to basename of folder_path)
    """
    if folder_name is None:
        folder_name = os.path.basename(folder_path)

    all_lines = []
    all_lines.append(f"# {folder_name} - Ultimate Note")
    all_lines.append("")
    all_lines.append("---")

    # PDFs
    for fpath in sorted(glob.glob(os.path.join(folder_path, "*.pdf"))):
        fname = os.path.basename(fpath)
        all_lines.append(f"\n## {fname}\n")
        all_lines.append(extract_pdf_text(fpath))
        all_lines.append("---")

    # PPTXs
    for fpath in sorted(glob.glob(os.path.join(folder_path, "*.pptx"))):
        fname = os.path.basename(fpath)
        all_lines.append(f"\n## {fname}\n")
        all_lines.append(extract_pptx_text(fpath))
        all_lines.append("---")

    # Binary PPTs (limited extraction)
    for fpath in sorted(glob.glob(os.path.join(folder_path, "*.ppt"))):
        fname = os.path.basename(fpath)
        all_lines.append(f"\n## {fname}\n")
        text = extract_ppt_clean(fpath)
        if text:
            all_lines.append(
                "*[Extracted from old .PPT format - text may be incomplete. "
                "Open the .ppt file to see slides/images.]*\n"
            )
            all_lines.append(text)
        else:
            all_lines.append(
                "*[Unable to extract text from old .PPT format. "
                "Open in PowerPoint to view.]*"
            )
        all_lines.append("\n---")

    # Markdown files (skip any existing ultimate note to avoid duplication)
    for fpath in sorted(glob.glob(os.path.join(folder_path, "*.md"))):
        fname = os.path.basename(fpath)
        if "ultimate" in fname.lower():
            continue
        all_lines.append(f"\n## {fname}\n")
        try:
            with open(fpath) as f:
                all_lines.append(f.read())
        except Exception as e:
            all_lines.append(f"*Error reading file: {e}*")
        all_lines.append("")

    output_path = os.path.join(folder_path, f"{folder_name} ultimate note.md")
    content = "\n".join(all_lines)
    with open(output_path, "w") as f:
        f.write(content)

    print(f"✓ {output_path} ({len(all_lines)} lines, {os.path.getsize(output_path)} bytes)")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Compile an 'ultimate note' markdown from course materials."
    )
    parser.add_argument("folder", help="Path to the module folder containing PDFs, PPTXs, etc.")
    parser.add_argument("--name", help="Display name for the module (defaults to folder basename)")
    args = parser.parse_args()
    compile_ultimate_note(args.folder, args.name)
